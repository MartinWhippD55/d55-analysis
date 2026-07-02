"""
Generate BRYT Contract Note Rework presentation using D55 template.
Creates an exec-level deck walking through the 5 estimates with high-level figures.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy

# Load D55 template
template_path = 'analysis/D55/ai-dlc/assets/powerpoint/D55_Deck_Visual_NO ANIMATION Template (1).pptx'
prs = Presentation(template_path)

# Print available layouts for reference
print("Available slide layouts:")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  {i}: {layout.name}")

# Remove all existing slides (we'll build from scratch using the layouts)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Helper functions
def add_slide(layout_idx):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def clear_placeholders(slide):
    for ph in slide.placeholders:
        ph.text_frame.clear()

def set_text(text_frame, text, size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = add_slide(0)  # 1_Title layout
clear_placeholders(slide)
# Use placeholders from title layout
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:  # Title
        set_text(ph.text_frame, "BRYT Energy\nContract Note Rework", size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    elif ph.placeholder_format.idx == 1:  # Subtitle area
        set_text(ph.text_frame, "Estimate Overview", size=16, color=RGBColor(0xFF, 0xFF, 0xFF))
    elif "date" in str(ph.placeholder_format.idx).lower() or ph.placeholder_format.idx == 10:
        set_text(ph.text_frame, "July 2026", size=12, color=RGBColor(0xFF, 0xFF, 0xFF))

# ============================================================
# SLIDE 2: Overview / Contents
# ============================================================
slide = add_slide(4)  # 1_Content layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Overview", size=24, bold=True)
    elif ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        lines = [
            "Rework of the BRYT contract note system across 5 estimates:",
            "",
            "1. PDF / Template Management — Self-service template editor",
            "2. DocuSign Integration — Automated e-signature flow",
            "3. Training & Data Sources — Enablement + extensibility",
            "4. Bespoke Contracts — Custom one-off contract notes",
            "5. Comparison Audit — Detect PDF tampering before sending",
            "",
            "Total: ~58 developer days (45 required + 13 optional)"
        ]
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(14)
            if line.startswith("Total:"):
                run.font.bold = True

# ============================================================
# SLIDE 3: Summary Table
# ============================================================
slide = add_slide(18)  # Table SLIDE layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Estimate Summary", size=24, bold=True)

# Add table
from pptx.util import Inches
rows, cols = 8, 4
left = Inches(0.8)
top = Inches(2.0)
width = Inches(8.4)
height = Inches(3.5)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(1.8)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(1.6)

# Headers
headers = ['Estimate', 'Required (days)', 'Optional', 'Total']
for col_idx, h in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)

# Data rows
data = [
    ('1. PDF / Template Management', '9.0', '4.5', '13.5'),
    ('2. DocuSign Integration', '4.1', '3.0', '7.1'),
    ('3a. Training & Enablement', '8.0', '0', '8.0'),
    ('3b. Data Source Extensibility', '6.4', '2.5', '8.9'),
    ('4. Bespoke Contracts', '4.8', '3.0', '7.8'),
    ('5. Comparison Audit', '12.4', '0', '12.4'),
    ('TOTAL', '44.6', '13.0', '57.6'),
]

for row_idx, row_data in enumerate(data, 1):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            if row_idx == len(data):
                p.font.bold = True

# ============================================================
# SLIDE 4: Est 1 - PDF/Template Management
# ============================================================
slide = add_slide(7)  # 1_Section Slide
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "01", size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    elif ph.placeholder_format.idx == 1:
        set_text(ph.text_frame, "PDF / Template\nManagement", size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# ============================================================
# SLIDE 5: Est 1 Detail
# ============================================================
slide = add_slide(4)  # Content layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Est 1: PDF / Template Management", size=22, bold=True)
    elif ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        lines = [
            "Replace the current developer-dependent SVG/HTML pipeline with a",
            "visual, self-service template management system.",
            "",
            "Key deliverables:",
            "• Template CRUD with priority-based selection rules",
            "• Visual section editor (pdf-me Designer embedded in Angular)",
            "• Shared/reusable sections (headers, footers, T&Cs)",
            "• Rules engine for automated template matching",
            "• Render pipeline (Lambda: section render + PDF stitch)",
            "• Version history with revert capability",
            "",
            "Estimate: 13.5 days (9.0 required + 4.5 optional testing)",
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            if line.startswith("Estimate:"):
                run.font.bold = True

# ============================================================
# SLIDE 6: Est 2 - DocuSign
# ============================================================
slide = add_slide(7)  # Section Slide
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "02", size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    elif ph.placeholder_format.idx == 1:
        set_text(ph.text_frame, "DocuSign\nIntegration", size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# ============================================================
# SLIDE 7: Est 2 Detail
# ============================================================
slide = add_slide(4)  # Content layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Est 2: DocuSign Integration", size=22, bold=True)
    elif ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        lines = [
            "Automated e-signature pipeline: PDF rendered → sent for signing →",
            "signed copy stored back in Salesforce.",
            "",
            "Key deliverables:",
            "• Automated S3 trigger (PDF created → DocuSign flow starts)",
            "• Salesforce customer lookup (contact details for signing)",
            "• DocuSign JWT auth + envelope creation",
            "• Webhook handler (signed PDF → S3 + Salesforce attachment)",
            "• DynamoDB envelope tracking for debugging",
            "",
            "Estimate: 7.1 days (4.1 required + 3.0 optional testing)",
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            if line.startswith("Estimate:"):
                run.font.bold = True

# ============================================================
# SLIDE 8: Est 3 - Training & Data Sources
# ============================================================
slide = add_slide(7)  # Section Slide
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "03", size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    elif ph.placeholder_format.idx == 1:
        set_text(ph.text_frame, "Training &\nData Sources", size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# ============================================================
# SLIDE 9: Est 3 Detail
# ============================================================
slide = add_slide(4)  # Content layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Est 3: Training & Data Sources", size=22, bold=True)
    elif ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        lines = [
            "Two sub-deliverables: enablement assets + data source extensibility.",
            "",
            "3a. Training & Enablement (8 days):",
            "• Quick-start guide, how-to guides, data field reference",
            "• Rules engine cheat sheet, troubleshooting guide",
            "• Screen recordings (optional, post-build)",
            "",
            "3b. Data Source Extensibility (8.9 days):",
            "• Glue Data Catalog discovery via SageMaker Unified Studio",
            "• Template-level data source attachment",
            "• Athena enrichment at render time (keyed on BrytNumber)",
            "• Field browser in Section Editor",
            "",
            "Combined estimate: 16.9 days",
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            if line.startswith("Combined"):
                run.font.bold = True

# ============================================================
# SLIDE 10: Est 4 - Bespoke Contracts
# ============================================================
slide = add_slide(7)  # Section Slide
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "04", size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    elif ph.placeholder_format.idx == 1:
        set_text(ph.text_frame, "Bespoke\nContracts", size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# ============================================================
# SLIDE 11: Est 4 Detail
# ============================================================
slide = add_slide(4)  # Content layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Est 4: Bespoke Contracts", size=22, bold=True)
    elif ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        lines = [
            "Custom contract notes for VIP customers or non-standard terms.",
            "Pipeline skips these; users create them manually in the portal.",
            "",
            "Key deliverables:",
            "• Pipeline skip for bespoke-flagged customers (Salesforce field)",
            "• Bespoke contract list (pending, draft, rendered statuses)",
            "• Editor with clone-from-template or start-from-scratch",
            "• On-demand render (Save & Render button)",
            "• Manual DocuSign trigger (Send via DocuSign button)",
            "• Contract data reference panel + version history",
            "",
            "Estimate: 7.8 days (4.8 required + 3.0 optional testing)",
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            if line.startswith("Estimate:"):
                run.font.bold = True

# ============================================================
# SLIDE 12: Est 5 - Comparison Audit
# ============================================================
slide = add_slide(7)  # Section Slide
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "05", size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    elif ph.placeholder_format.idx == 1:
        set_text(ph.text_frame, "Comparison\nAudit", size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# ============================================================
# SLIDE 13: Est 5 Detail
# ============================================================
slide = add_slide(4)  # Content layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Est 5: Comparison Audit", size=22, bold=True)
    elif ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        lines = [
            "Detect whether contract note PDFs are being manually edited",
            "between rendering and sending to customers.",
            "",
            "Key deliverables:",
            "• Step Function pipeline (batch processing)",
            "• Microsoft Graph API integration (fetch sent PDFs from Outlook)",
            "• AWS Bedrock comparison (AI-powered diff analysis)",
            "• S3 results storage (Athena-queryable, date-partitioned)",
            "• Prompt iteration (3-5 cycles to refine output quality)",
            "• Spreadsheet reporting for BRYT",
            "",
            "Estimate: 12.4 days (includes prompt iteration budget)",
            "",
            "Dependency: Requires Graph API access (M365 admin approval)",
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            if line.startswith("Estimate:") or line.startswith("Dependency:"):
                run.font.bold = True

# ============================================================
# SLIDE 14: Open Questions
# ============================================================
slide = add_slide(4)  # Content layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Open Questions", size=22, bold=True)
    elif ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        lines = [
            "Key items requiring client confirmation before/during implementation:",
            "",
            "• Multi-party signing (customer only, or BRYT rep + TPI too)?",
            "• DocuSign account — new setup required (no existing account found)",
            "• Email branding — DocuSign standard emails acceptable?",
            "• Admin Portal status UI — needed, or backend-only?",
            "• Salesforce object mapping for customersalesforceref",
            "• Voided envelopes / resend / reminders — out of scope?",
            "• DocuSign for bespoke contracts — manual trigger acceptable?",
            "• Bespoke flag mechanism in Salesforce",
            "• Graph API access for Outlook mail search (M365 admin)",
            "• Mailbox identification + email correlation method",
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(12)

# ============================================================
# SLIDE 15: Next Steps
# ============================================================
slide = add_slide(4)  # Content layout
clear_placeholders(slide)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        set_text(ph.text_frame, "Next Steps", size=22, bold=True)
    elif ph.placeholder_format.idx == 1:
        tf = ph.text_frame
        tf.clear()
        lines = [
            "1. Resolve open questions (client confirmation needed)",
            "2. Prioritise estimates (sequential or parallel delivery)",
            "3. Confirm optional scope (testing tasks — recommended)",
            "4. Begin implementation from Estimate 1",
            "",
            "Estimates are designed to be delivered sequentially,",
            "each building on the previous. However, Est 5 (Comparison",
            "Audit) can run independently at any time.",
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(14)

# Save
output_path = 'analysis/BRYT/contract-note/BRYT Contract Note Rework - Estimates.pptx'
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
