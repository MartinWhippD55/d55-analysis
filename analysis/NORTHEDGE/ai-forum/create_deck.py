from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DECK_DIR = os.path.dirname(os.path.abspath(__file__))
TITLE_BG = os.path.join(DECK_DIR, "D55.title.slide.png")
CONTENT_BG = os.path.join(DECK_DIR, "D55.background.slide.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colours
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x00, 0xE5, 0xD0)
GREY = RGBColor(0xAA, 0xAA, 0xAA)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)

def add_bg_image(slide, img_path):
    """Add a full-slide background image."""
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_bullet_slide(slide, title, bullets, subtitle=None):
    add_bg_image(slide, CONTENT_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 title, font_size=36, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.8),
                     subtitle, font_size=18, bold=False, color=CYAN)
    start_y = Inches(2.2) if subtitle else Inches(1.8)
    txBox = slide.shapes.add_textbox(Inches(0.8), start_y, Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

# --- SLIDE 1: Title ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(slide, TITLE_BG)
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.5),
             "AI-First Delivery", font_size=48, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(1),
             "Practical lessons from the field", font_size=24, bold=False, color=LIGHT_GREY)
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(0.6),
             "D55  |  Northedge AI Forum", font_size=16, bold=False, color=GREY)

# --- SLIDE 2: Agenda ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Agenda", [
    "1.  People & Adoption — why transformations fail",
    "2.  How AI-First Changes Delivery — the new lifecycle",
    "3.  The Economics — $200/seat, 5x output",
    "4.  First 8 Weeks — how we'd start",
    "5.  What \"Good\" Looks Like — for a PE exit"
])

# --- SLIDE 3: People - The Problem ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "People & Adoption", [
    "The tech is the easy part.",
    "Getting people to change is where AI transformations die.",
    "",
    "On our recent project (10 devs, AI-first):",
    "•  Primary resistance came from POs, BAs, Programme Managers",
    "•  When devs go 5x faster, coordination roles become the bottleneck",
    "•  The \"middle layer\" correctly senses their role is changing"
], subtitle="The real challenge isn't technology")

# --- SLIDE 4: Four Fears Framework ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(slide, CONTENT_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Four Fears Framework", font_size=36, bold=True, color=WHITE)

# Axis labels
add_text_box(slide, Inches(3.0), Inches(1.2), Inches(3), Inches(0.5),
             "Past / Present", font_size=14, color=GREY, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.5), Inches(1.2), Inches(3), Inches(0.5),
             "Future Thinking", font_size=14, color=GREY, align=PP_ALIGN.CENTER)

box_w = Inches(5.2)
box_h = Inches(2.5)
left1 = Inches(1.0)
left2 = Inches(6.8)
top1 = Inches(1.7)
top2 = Inches(4.5)

# Blue - Loss of Control
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left1, top1, box_w, box_h)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
shape.line.color.rgb = RGBColor(0x42, 0xA5, 0xF5)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Loss of Control"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(0x42, 0xA5, 0xF5)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Controller"
p2.font.size = Pt(14)
p2.font.color.rgb = LIGHT_GREY
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "\n↻ More oversight & observability, not less"
p3.font.size = Pt(14)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER

# Red - Project Failure
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left2, top1, box_w, box_h)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x5C, 0x1A, 0x1A)
shape.line.color.rgb = RGBColor(0xEF, 0x53, 0x50)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Project Failure"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(0xEF, 0x53, 0x50)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Driver"
p2.font.size = Pt(14)
p2.font.color.rgb = LIGHT_GREY
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "\n↻ Faster, cheaper failure — compressed risk"
p3.font.size = Pt(14)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER

# Green - Disruption to Hierarchy
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left1, top2, box_w, box_h)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x4C, 0x1A)
shape.line.color.rgb = RGBColor(0x66, 0xBB, 0x6A)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Disruption to Hierarchy"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(0x66, 0xBB, 0x6A)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Stabiliser"
p2.font.size = Pt(14)
p2.font.color.rgb = LIGHT_GREY
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "\n↻ Orchestrators become more valuable"
p3.font.size = Pt(14)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER

# Yellow - Feeling Overshadowed
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left2, top2, box_w, box_h)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x4C, 0x3D, 0x0A)
shape.line.color.rgb = RGBColor(0xFF, 0xCA, 0x28)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Feeling Overshadowed"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xCA, 0x28)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Influencer"
p2.font.size = Pt(14)
p2.font.color.rgb = LIGHT_GREY
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "\n↻ Amplified, not replaced — more stage"
p3.font.size = Pt(14)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER

# --- SLIDE 5: What unblocked ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "What Unblocked Adoption", [
    "•  Action first, permission later — prove it, then present results",
    "•  Metrics from day one — 5x improvement is hard to argue with",
    "•  Regular retros — safe space to surface tension early",
    "•  1-on-1 training & support — catch people struggling quietly",
    "•  Accept not everyone will come around — plan for it"
], subtitle="From our recent AI-first project")

# --- SLIDE 6: AI-First Delivery ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "How AI-First Changes Delivery", [
    "1.  Analysis & Design — more time here, produce clear specifications",
    "2.  Contracts & Integration — understand touch points upfront",
    "3.  Implementation — devs with agentic tooling against defined specs",
    "4.  Testing — human-in-the-loop + AI for unit tests & seed data",
    "5.  Documentation — AI-assisted handover assets",
    "",
    "Key insight: when implementation is 5x faster, the bottleneck",
    "shifts to clarity of intent. Design discipline is the multiplier",
    "on top of the multiplier."
], subtitle="Not waterfall — design-led agile at speed")

# --- SLIDE 7: Economics ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(slide, CONTENT_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
             "The Economics", font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.5),
             "$200 / seat / month", font_size=56, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(1),
             "5x developer productivity", font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
bullets = [
    "Copilot + Claude seats = ~$200/dev/month",
    "For 100 devs: $20k/month for 5x output",
    "Equivalent of 400+ additional developers",
    "Ongoing cost is primarily tooling — no custom infrastructure",
    "Even 20% headcount reduction = millions saved annually"
]
for i, b in enumerate(bullets):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = "•  " + b
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.space_after = Pt(8)

# --- SLIDE 8: First 8 Weeks ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(slide, CONTENT_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
             "First 8 Weeks", font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Assess. Prove. Translate.", font_size=20, bold=False, color=CYAN)

col_w = Inches(3.5)
col_top = Inches(2.2)

# Week 1-2
txBox = slide.shapes.add_textbox(Inches(0.8), col_top, col_w, Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Weeks 1–2: Assess"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = CYAN
for b in ["Assess codebase & tooling", "Map current AI adoption", "Identify quick wins"]:
    p = tf.add_paragraph()
    p.text = "•  " + b
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.space_after = Pt(6)

# Week 3-6
txBox = slide.shapes.add_textbox(Inches(4.8), col_top, col_w, Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Weeks 3–6: Prove"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = CYAN
for b in ["Embed D55 or stand up squad", "AI-first on real work", "Track metrics from day 1"]:
    p = tf.add_paragraph()
    p.text = "•  " + b
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.space_after = Pt(6)

# Week 7-8
txBox = slide.shapes.add_textbox(Inches(8.8), col_top, col_w, Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Weeks 7–8: Translate"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = CYAN
for b in ["Demonstrate with numbers", "Produce reusable playbook", "Rollout plan for other teams"]:
    p = tf.add_paragraph()
    p.text = "•  " + b
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.space_after = Pt(6)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             "D55 embeds, proves, then enables. Not a permanent dependency.",
             font_size=16, bold=False, color=GREY)

# --- SLIDE 9: What Good Looks Like ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "What \"Good\" Looks Like", [
    "•  100 devs producing the output of 300–500",
    "•  Release cadence in days, not months",
    "•  AI embedded in the product (compliance, document gen, self-service)",
    "•  Engineering org that acquirers see as an asset, not a cost",
    "",
    "What to avoid:",
    "•  Don't boil the ocean — pick 2–3 squads, prove it, roll out",
    "•  Don't start with AI product features before engineering is AI-native",
    "•  This is a people & delivery transformation — not a tech project"
], subtitle="PE exit positioning")

# --- SLIDE 10: Close ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_image(slide, TITLE_BG)
add_text_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(1.5),
             "We've done this.", font_size=44, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(1),
             "Let us show you how.", font_size=28, bold=False, color=LIGHT_GREY)

# Save
output = os.path.join(DECK_DIR, "D55-AI-Forum-Deck.pptx")
prs.save(output)
print(f"Done — saved to {output}")
